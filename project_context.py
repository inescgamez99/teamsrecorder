"""
Memoria de proyecto híbrida (estilo Claude Code).

- Python extrae los documentos vinculados (Word/PPT/PDF/Excel/txt/md/csv) a ficheros
  .txt dentro de `project_docs/<project_id>/docs/` (incremental: solo re-extrae los
  que cambian).
- Los resúmenes de reuniones pasadas del proyecto se guardan en
  `project_docs/<project_id>/meetings/` como .txt.
- Al generar minutas se le da a `claude -p` acceso agéntico a esa carpeta
  (`--add-dir`) para que busque y lea lo relevante por sí mismo. Escala a 100+ docs
  porque Claude solo lee lo que necesita para esa reunión.
"""
import json
import logging
import re
from pathlib import Path

from config import PROJECT_DIR, MINUTES_DIR

log = logging.getLogger(__name__)

_MAX_PER_FILE = 20000   # chars por documento extraído (Claude ya elige qué leer)
_SUPPORTED = {'.txt', '.md', '.csv', '.docx', '.pptx', '.pdf', '.xlsx'}


def project_docs_dir(project_id: str) -> Path:
    return PROJECT_DIR / 'project_docs' / project_id


def load_projects() -> list:
    p = PROJECT_DIR / 'projects.json'
    if p.exists():
        try:
            return json.loads(p.read_text(encoding='utf-8')).get('projects', [])
        except Exception:
            pass
    return []


def _norm(s: str) -> str:
    return re.sub(r'\s+', ' ', re.sub(r'[^\w\s]', ' ', (s or '').lower())).strip()


def detect_project(text: str, meeting_name: str = '', projects: list = None) -> dict | None:
    """Empareja la reunión con un proyecto por nombre/id (palabras clave que
    aparezcan en el nombre de la reunión o en la transcripción)."""
    projects = projects if projects is not None else load_projects()
    hay = _norm((meeting_name or '') + ' ' + (text or '')[:6000])
    if not hay:
        return None
    for p in projects:
        name = _norm(p.get('name', ''))
        pid = _norm(p.get('id', ''))
        keys = [k for k in (name, pid) if k]
        words = [w for w in name.split() if len(w) > 3]
        if any(k in hay for k in keys) or any(w in hay for w in words):
            return p
    return None


# ── Extracción de texto por tipo de archivo ──────────────────────────────────
def _extract_file_text(path: Path) -> str:
    ext = path.suffix.lower()
    try:
        if ext in ('.txt', '.md', '.csv'):
            return path.read_text(encoding='utf-8', errors='ignore')[:_MAX_PER_FILE]
        if ext == '.docx':
            import docx
            doc = docx.Document(str(path))
            return '\n'.join(p.text for p in doc.paragraphs if p.text.strip())[:_MAX_PER_FILE]
        if ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(str(path))
            out = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if getattr(shape, 'has_text_frame', False) and shape.text_frame.text.strip():
                        out.append(shape.text_frame.text.strip())
            return '\n'.join(out)[:_MAX_PER_FILE]
        if ext == '.pdf':
            from pypdf import PdfReader
            r = PdfReader(str(path))
            out = []
            for page in r.pages[:40]:
                out.append(page.extract_text() or '')
                if sum(len(x) for x in out) > _MAX_PER_FILE:
                    break
            return '\n'.join(out)[:_MAX_PER_FILE]
        if ext == '.xlsx':
            import openpyxl
            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            out = []
            for ws in wb.worksheets[:8]:
                out.append(f"[Hoja: {ws.title}]")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        out.append(' | '.join(cells))
                    if sum(len(x) for x in out) > _MAX_PER_FILE:
                        break
                if sum(len(x) for x in out) > _MAX_PER_FILE:
                    break
            wb.close()
            return '\n'.join(out)[:_MAX_PER_FILE]
    except Exception as e:
        log.warning(f"project_context extract {path.name}: {e}")
    return ''


def _safe_name(p: Path) -> str:
    return re.sub(r'[^\w.\- ]', '_', p.stem)[:90].strip() + '.txt'


_MAX_FILE_SIZE = 40 * 1024 * 1024   # 40 MB — saltar ficheros gigantes
_MAX_FILES = 1500                    # tope de ficheros a considerar


def _collect_files(dirs, should_cancel=None) -> list:
    """Recopila los ficheros soportados de las carpetas (incluye los de OneDrive
    'solo en la nube', que se descargarán al leerlos). Filtra solo por tamaño y
    número para mantener el control."""
    files = []
    for d in (dirs or []):
        folder = Path(d)
        if not folder.exists() or not folder.is_dir():
            continue
        try:
            for f in folder.rglob('*'):
                if should_cancel and should_cancel():
                    return files
                if (not f.is_file() or f.suffix.lower() not in _SUPPORTED
                        or f.name.startswith('~$')):
                    continue
                try:
                    if f.stat().st_size > _MAX_FILE_SIZE:
                        continue
                except Exception:
                    continue
                files.append(f)
                if len(files) >= _MAX_FILES:
                    return files
        except Exception as e:
            log.warning(f"_collect_files {d}: {e}")
    return files


def _unique_out_name(f: Path, taken: set) -> str:
    base = _safe_name(f)
    if base not in taken:
        return base
    stem = base[:-4]  # sin '.txt'
    i = 2
    while f"{stem}_{i}.txt" in taken:
        i += 1
    return f"{stem}_{i}.txt"


def sync_project_docs(project: dict, progress_cb=None, should_cancel=None) -> int:
    """Sincroniza la memoria de documentos con las carpetas vinculadas (espejo real):
    extrae los documentos nuevos/cambiados y ELIMINA el texto de los que ya no existen
    (carpetas o ficheros quitados). Incremental; con progreso y cancelable.
    Devuelve el nº de documentos (re)extraídos."""
    pid = project.get('id')
    dirs = project.get('context_dirs', []) or []
    if not pid:
        return 0
    out_dir = project_docs_dir(pid) / 'docs'
    out_dir.mkdir(parents=True, exist_ok=True)
    manifest_path = project_docs_dir(pid) / '.sync.json'
    manifest = {}
    if manifest_path.exists():
        try:
            raw = json.loads(manifest_path.read_text(encoding='utf-8'))
            # formato nuevo: {source: {mtime, out}}; tolera formato viejo (float)
            for k, v in raw.items():
                manifest[k] = v if isinstance(v, dict) else {'mtime': v, 'out': None}
        except Exception:
            pass

    files = _collect_files(dirs, should_cancel)
    current_keys = {str(f) for f in files}
    taken = {v.get('out') for v in manifest.values() if v.get('out')}

    total = len(files)
    synced = 0
    for i, f in enumerate(files):
        if should_cancel and should_cancel():
            break
        if progress_cb:
            progress_cb(i + 1, total, f.name)
        key = str(f)
        try:
            mtime = f.stat().st_mtime
        except Exception:
            continue
        prev = manifest.get(key)
        out_name = prev.get('out') if (prev and prev.get('out')) else _unique_out_name(f, taken)
        taken.add(out_name)
        if prev and prev.get('mtime') == mtime and (out_dir / out_name).exists():
            continue  # sin cambios
        txt = _extract_file_text(f)
        if txt.strip():
            header = f"# Documento: {f.name}\n# Ruta original: {f}\n\n"
            try:
                (out_dir / out_name).write_text(header + txt, encoding='utf-8')
                manifest[key] = {'mtime': mtime, 'out': out_name}
                synced += 1
            except Exception as e:
                log.warning(f"sync write {f.name}: {e}")

    # Espejo: eliminar los documentos cuya fuente ya no existe (solo si NO se canceló,
    # porque al cancelar la lista de ficheros podría estar incompleta).
    pruned = 0
    if not (should_cancel and should_cancel()):
        for k in [k for k in manifest if k not in current_keys]:
            out = manifest[k].get('out')
            if out:
                try:
                    (out_dir / out).unlink(missing_ok=True)
                except Exception:
                    pass
            del manifest[k]
            pruned += 1

    try:
        manifest_path.write_text(json.dumps(manifest), encoding='utf-8')
    except Exception:
        pass
    if synced or pruned:
        log.info(f"sync_project_docs '{pid}': {synced}/{total} extraídos, {pruned} eliminados")
    return synced


def _extract_summary(md_text: str) -> str:
    m = re.search(
        r'##\s*(Resumen Ejecutivo|Executive Summary)\b(.*?)(?=\n##\s|\Z)',
        md_text, re.IGNORECASE | re.DOTALL,
    )
    if m:
        body = m.group(2).strip()
        # incluir también Decisiones si están
        d = re.search(r'##\s*(Decisiones Tomadas|Decisions Made)\b(.*?)(?=\n##\s|\Z)',
                      md_text, re.IGNORECASE | re.DOTALL)
        if d:
            body += "\n\n## " + d.group(1) + d.group(2).rstrip()
        return body[:2500]
    return md_text.strip()[:1500]


def add_meeting_summary(project_id: str, stem: str, title: str, date: str, minutes_text: str) -> None:
    """Guarda el resumen de una reunión como contexto para futuras reuniones del proyecto."""
    if not project_id or project_id == 'none':
        return
    try:
        mdir = project_docs_dir(project_id) / 'meetings'
        mdir.mkdir(parents=True, exist_ok=True)
        summary = _extract_summary(minutes_text)
        (mdir / f"{stem}.txt").write_text(
            f"# Reunión: {title} ({date})\n\n{summary}\n", encoding='utf-8')
        log.info(f"add_meeting_summary '{project_id}': {title}")
    except Exception as e:
        log.warning(f"add_meeting_summary: {e}")


def get_context_dir(project: dict) -> str | None:
    """Devuelve la ruta de project_docs/<id> si tiene contenido (para --add-dir)."""
    pid = project.get('id') if project else None
    if not pid:
        return None
    pdir = project_docs_dir(pid)
    for sub in ('docs', 'meetings'):
        p = pdir / sub
        if p.exists() and any(p.iterdir()):
            return str(pdir)
    return None


def prepare_context(transcript: str, meeting_name: str = '') -> tuple:
    """Detecta el proyecto, sincroniza sus documentos y devuelve
    (project_dict|None, context_dir|None) para pasar a la generación de minutas."""
    proj = detect_project(transcript, meeting_name)
    if not proj:
        return None, None
    try:
        sync_project_docs(proj)
    except Exception as e:
        log.warning(f"prepare_context sync: {e}")
    return proj, get_context_dir(proj)
